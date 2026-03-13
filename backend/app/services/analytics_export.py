"""analytics_export.py — Dashboard PPTX export service.

Converts an Analytics dashboard (list of widgets with titles + SQL results)
into a PowerPoint presentation using python-pptx (already in requirements.txt).

One slide per widget:
  - Title slide at position 0 with the dashboard name
  - Table slides for row-based data
  - Numeric/KPI slides for single-value metrics
  - Placeholder slides for chart types (chart image would require Playwright —
    falls back to a table of the underlying data)
"""
from __future__ import annotations

import io
from datetime import datetime
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Urban Vibes Dynamics brand colours
_PRIMARY = RGBColor(0x51, 0x45, 0x9D)   # #51459d
_SUCCESS = RGBColor(0x6F, 0xD9, 0x43)   # #6fd943
_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
_GRAY    = RGBColor(0xF3, 0xF4, 0xF6)
_TEXT    = RGBColor(0x1F, 0x29, 0x37)

# Slide dimensions: widescreen 16:9
_SLIDE_W = Inches(13.33)
_SLIDE_H = Inches(7.5)


def _set_bg(slide, color: RGBColor) -> None:
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _title_slide(prs: Presentation, dashboard_name: str, exported_at: str) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, _PRIMARY)

    # Title text box
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = dashboard_name
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = _WHITE

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.33), Inches(0.6))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = f"Exported {exported_at} · Urban Vibes Dynamics Analytics"
    run2.font.size = Pt(14)
    run2.font.color.rgb = RGBColor(0xC7, 0xC2, 0xE8)


def _widget_slide(
    prs: Presentation,
    title: str,
    widget_type: str,
    rows: list[dict[str, Any]],
) -> None:
    """Add a slide for a single widget."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, _WHITE)

    # Header bar
    from pptx.oxml.ns import qn
    header = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, _SLIDE_W, Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = _PRIMARY
    header.line.fill.background()

    # Widget title
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = _WHITE

    if not rows:
        _add_no_data_text(slide)
        return

    # KPI / single-value widget
    if len(rows) == 1 and len(rows[0]) == 1:
        _kpi_content(slide, rows[0])
        return

    # Table widget (default)
    _table_content(slide, rows)


def _add_no_data_text(slide) -> None:
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "No data available"
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)


def _kpi_content(slide, row: dict[str, Any]) -> None:
    key = list(row.keys())[0]
    value = list(row.values())[0]
    txBox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9.33), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(value)
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = _PRIMARY

    # Label below
    txLabel = slide.shapes.add_textbox(Inches(2), Inches(4.7), Inches(9.33), Inches(0.5))
    tf2 = txLabel.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = key
    r2.font.size = Pt(14)
    r2.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)


def _table_content(slide, rows: list[dict[str, Any]]) -> None:
    if not rows:
        return
    cols = list(rows[0].keys())
    n_cols = len(cols)
    max_rows = min(len(rows), 15)  # cap at 15 rows per slide

    col_width = Inches(12.5) / n_cols
    row_height = Inches(0.38)
    top = Inches(1.1)
    left = Inches(0.4)

    table = slide.shapes.add_table(
        max_rows + 1, n_cols,
        left, top,
        Inches(12.5), row_height * (max_rows + 1)
    ).table

    # Header row
    for ci, col_name in enumerate(cols):
        cell = table.cell(0, ci)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0] if p.runs else p.add_run()
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = _WHITE

    # Data rows
    for ri, row in enumerate(rows[:max_rows]):
        for ci, col_name in enumerate(cols):
            cell = table.cell(ri + 1, ci)
            cell.text = str(row.get(col_name, ""))
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _GRAY
            p = cell.text_frame.paragraphs[0]
            r = p.runs[0] if p.runs else p.add_run()
            r.font.size = Pt(10)
            r.font.color.rgb = _TEXT


def export_dashboard_to_pptx(
    dashboard_name: str,
    widgets: list[dict[str, Any]],
) -> bytes:
    """Convert a dashboard to a PPTX file and return the bytes.

    Args:
        dashboard_name: Display name shown on the title slide.
        widgets: List of widget dicts, each with keys:
                   - title (str)
                   - type  (str): 'number'|'bar'|'line'|'pie'|'table'|etc.
                   - rows  (list[dict]): query result rows

    Returns:
        Raw bytes of the .pptx file.
    """
    prs = Presentation()
    prs.slide_width  = _SLIDE_W
    prs.slide_height = _SLIDE_H

    exported_at = datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC")

    # Title slide
    _title_slide(prs, dashboard_name, exported_at)

    # One slide per widget
    for widget in widgets:
        _widget_slide(
            prs,
            title=widget.get("title", "Widget"),
            widget_type=widget.get("type", "table"),
            rows=widget.get("rows", []),
        )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
